"""
Tests for file extraction functionality.

Tests PDF, DOCX, and PPTX text extraction using real test files.
"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF

from main import (
    extract_pdf,
    extract_docx,
    extract_pptx,
    count_words,
    get_extension_from_content_type,
)


# Fixtures for creating test files
@pytest.fixture
def sample_pdf() -> Path:
    """Create a sample PDF file for testing."""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    temp_path = Path(temp_file.name)
    temp_file.close()
    
    doc = fitz.open()
    
    # Page 1
    page1 = doc.new_page()
    page1.insert_text((72, 72), "This is the first page of the PDF document.")
    page1.insert_text((72, 100), "It contains sample text for testing extraction.")
    
    # Page 2
    page2 = doc.new_page()
    page2.insert_text((72, 72), "This is the second page.")
    page2.insert_text((72, 100), "More content here for testing purposes.")
    
    doc.save(temp_path)
    doc.close()
    
    yield temp_path
    
    # Cleanup
    if temp_path.exists():
        temp_path.unlink()


@pytest.fixture
def sample_docx() -> Path:
    """Create a sample DOCX file for testing."""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    temp_path = Path(temp_file.name)
    temp_file.close()
    
    doc = Document()
    doc.add_paragraph("This is the first paragraph of the document.")
    doc.add_paragraph("This is the second paragraph with more content.")
    doc.add_paragraph("Final paragraph for testing extraction.")
    
    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"
    
    doc.save(temp_path)
    
    yield temp_path
    
    # Cleanup
    if temp_path.exists():
        temp_path.unlink()


@pytest.fixture
def sample_pptx() -> Path:
    """Create a sample PPTX file for testing."""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    temp_path = Path(temp_file.name)
    temp_file.close()
    
    prs = Presentation()
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(1)
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Welcome to the presentation"
    
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    textbox2 = slide2.shapes.add_textbox(left, top, width, height)
    textbox2.text_frame.text = "This is slide two content"
    
    prs.save(temp_path)
    
    yield temp_path
    
    # Cleanup
    if temp_path.exists():
        temp_path.unlink()


# PDF Extraction Tests
class TestPDFExtraction:
    """Tests for PDF text extraction."""
    
    def test_extract_pdf_returns_text(self, sample_pdf: Path):
        """Test that PDF extraction returns text content."""
        text, page_count = extract_pdf(sample_pdf)
        
        assert text is not None
        assert len(text) > 0
        assert "first page" in text.lower()
    
    def test_extract_pdf_returns_correct_page_count(self, sample_pdf: Path):
        """Test that PDF extraction returns correct page count."""
        text, page_count = extract_pdf(sample_pdf)
        
        assert page_count == 2
    
    def test_extract_pdf_extracts_all_pages(self, sample_pdf: Path):
        """Test that PDF extraction includes content from all pages."""
        text, _ = extract_pdf(sample_pdf)
        
        assert "first page" in text.lower()
        assert "second page" in text.lower()


# DOCX Extraction Tests
class TestDOCXExtraction:
    """Tests for DOCX text extraction."""
    
    def test_extract_docx_returns_text(self, sample_docx: Path):
        """Test that DOCX extraction returns text content."""
        text, page_count = extract_docx(sample_docx)
        
        assert text is not None
        assert len(text) > 0
        assert "first paragraph" in text.lower()
    
    def test_extract_docx_extracts_all_paragraphs(self, sample_docx: Path):
        """Test that DOCX extraction includes all paragraphs."""
        text, _ = extract_docx(sample_docx)
        
        assert "first paragraph" in text.lower()
        assert "second paragraph" in text.lower()
        assert "final paragraph" in text.lower()
    
    def test_extract_docx_extracts_tables(self, sample_docx: Path):
        """Test that DOCX extraction includes table content."""
        text, _ = extract_docx(sample_docx)
        
        assert "header 1" in text.lower()
        assert "data 1" in text.lower()


# PPTX Extraction Tests
class TestPPTXExtraction:
    """Tests for PPTX text extraction."""
    
    def test_extract_pptx_returns_text(self, sample_pptx: Path):
        """Test that PPTX extraction returns text content."""
        text, slide_count = extract_pptx(sample_pptx)
        
        assert text is not None
        assert len(text) > 0
        assert "presentation" in text.lower()
    
    def test_extract_pptx_returns_correct_slide_count(self, sample_pptx: Path):
        """Test that PPTX extraction returns correct slide count."""
        text, slide_count = extract_pptx(sample_pptx)
        
        assert slide_count == 2
    
    def test_extract_pptx_extracts_all_slides(self, sample_pptx: Path):
        """Test that PPTX extraction includes content from all slides."""
        text, _ = extract_pptx(sample_pptx)
        
        assert "slide 1" in text.lower()
        assert "slide 2" in text.lower()


# Utility Function Tests
class TestUtilityFunctions:
    """Tests for utility functions."""
    
    def test_count_words(self):
        """Test word counting function."""
        assert count_words("hello world") == 2
        assert count_words("one two three four five") == 5
        assert count_words("") == 0
    
    def test_get_extension_from_content_type_pdf(self):
        """Test content type to extension mapping for PDF."""
        assert get_extension_from_content_type("application/pdf") == ".pdf"
    
    def test_get_extension_from_content_type_docx(self):
        """Test content type to extension mapping for DOCX."""
        content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        assert get_extension_from_content_type(content_type) == ".docx"
    
    def test_get_extension_from_content_type_pptx(self):
        """Test content type to extension mapping for PPTX."""
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert get_extension_from_content_type(content_type) == ".pptx"
    
    def test_get_extension_from_content_type_with_charset(self):
        """Test content type parsing with charset."""
        assert get_extension_from_content_type("application/pdf; charset=utf-8") == ".pdf"
    
    def test_get_extension_from_content_type_unknown(self):
        """Test unknown content type returns None."""
        assert get_extension_from_content_type("text/plain") is None
        assert get_extension_from_content_type(None) is None
