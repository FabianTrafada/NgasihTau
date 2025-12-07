"""
File Processor Service

A lightweight Python service for extracting text from PDF, DOCX, and PPTX files.
"""

import tempfile
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Literal

import httpx
from fastapi import FastAPI, HTTPException, status
from pydantic import BaseModel, Field, HttpUrl


@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan handler."""
    print("File Processor Service starting...")
    yield
    print("File Processor Service shutting down...")


app = FastAPI(
    title="File Processor Service",
    description="Extracts text content from PDF, DOCX, and PPTX files",
    version="0.1.0",
    lifespan=lifespan,
)


# Request/Response Models
class ExtractRequest(BaseModel):
    """Request model for text extraction."""

    file_url: HttpUrl = Field(..., description="URL to the file in MinIO/S3")
    file_type: Literal["pdf", "docx", "pptx"] = Field(
        ..., description="Type of file to extract"
    )


class ExtractMetadata(BaseModel):
    """Metadata about the extracted content."""

    pages: int = Field(default=0, description="Number of pages (for PDF/PPTX)")
    word_count: int = Field(default=0, description="Total word count")
    char_count: int = Field(default=0, description="Total character count")


class ExtractResponse(BaseModel):
    """Response model for text extraction."""

    text: str = Field(..., description="Extracted text content")
    metadata: ExtractMetadata = Field(..., description="Extraction metadata")


class HealthResponse(BaseModel):
    """Health check response."""

    status: str = Field(default="healthy")
    service: str = Field(default="file-processor")
    version: str = Field(default="1.0.0")
    timestamp: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())


class LiveResponse(BaseModel):
    """Liveness check response."""

    status: str = Field(default="alive")


class ReadyResponse(BaseModel):
    """Readiness check response."""

    status: str = Field(default="healthy")


class ErrorDetail(BaseModel):
    """Error response detail."""

    detail: str


# Text extraction functions
def extract_pdf(file_path: Path) -> tuple[str, int]:
    """
    Extract text from a PDF file using PyMuPDF.
    
    Returns:
        Tuple of (extracted_text, page_count)
    """
    import fitz  # PyMuPDF

    text_parts: list[str] = []
    page_count = 0

    with fitz.open(file_path) as doc:
        page_count = len(doc)
        for page in doc:
            page_text = page.get_text()
            if page_text:
                text_parts.append(page_text)

    return "\n\n".join(text_parts), page_count


def extract_docx(file_path: Path) -> tuple[str, int]:
    """
    Extract text from a DOCX file using python-docx.
    
    Returns:
        Tuple of (extracted_text, page_count)
        Note: DOCX doesn't have a reliable page count, returns 1
    """
    from docx import Document

    doc = Document(file_path)
    text_parts: list[str] = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_parts.append(" | ".join(row_text))

    # DOCX doesn't have reliable page count without rendering
    return "\n\n".join(text_parts), 1


def extract_pptx(file_path: Path) -> tuple[str, int]:
    """
    Extract text from a PPTX file using python-pptx.
    
    Returns:
        Tuple of (extracted_text, slide_count)
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts: list[str] = []
    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
            
            # Extract text from tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if slide_texts:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts), slide_count


def count_words(text: str) -> int:
    """Count words in text."""
    return len(text.split())


def get_extension_from_content_type(content_type: str | None) -> str | None:
    """
    Map content-type to file extension.
    
    Args:
        content_type: MIME type from response headers
    
    Returns:
        File extension (e.g., '.pdf') or None if unknown
    """
    content_type_map = {
        "application/pdf": ".pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        # Common variations
        "application/msword": ".doc",  # Old Word format (not supported but mapped)
        "application/vnd.ms-powerpoint": ".ppt",  # Old PowerPoint format
    }
    
    if content_type:
        # Handle content-type with charset (e.g., "application/pdf; charset=utf-8")
        base_type = content_type.split(";")[0].strip().lower()
        return content_type_map.get(base_type)
    
    return None


async def download_file(url: str, suffix: str) -> Path:
    """
    Download a file from MinIO/S3 URL to a temporary location.
    
    Args:
        url: URL to download from (MinIO presigned URL or direct URL)
        suffix: File extension (e.g., '.pdf')
    
    Returns:
        Path to the downloaded temporary file
    
    Raises:
        httpx.HTTPStatusError: If download fails with HTTP error
        httpx.RequestError: If network/connection error occurs
    """
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        response = await client.get(url)
        response.raise_for_status()
        
        # Optionally detect extension from content-type if available
        content_type = response.headers.get("content-type")
        detected_ext = get_extension_from_content_type(content_type)
        
        # Use detected extension if it matches expected type, otherwise use provided suffix
        # This provides a safety check that the file type matches expectations
        final_suffix = suffix
        if detected_ext and detected_ext != suffix:
            # Log warning but use the provided suffix (trust the caller)
            print(f"Warning: Content-Type suggests {detected_ext} but expected {suffix}")
        
        # Create temp file with proper extension
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=final_suffix)
        temp_file.write(response.content)
        temp_file.close()
        
        return Path(temp_file.name)


# Health Endpoints
@app.get("/health", response_model=HealthResponse, tags=["Health"])
async def health_check() -> HealthResponse:
    """
    Detailed health check endpoint.

    Returns the current health status of the service with metadata.
    """
    return HealthResponse()


@app.get("/health/live", response_model=LiveResponse, tags=["Health"])
async def liveness_check() -> LiveResponse:
    """
    Liveness probe endpoint.

    Simple check that returns if the service is alive.
    Used by Kubernetes liveness probes.
    """
    return LiveResponse()


@app.get("/ready", response_model=ReadyResponse, tags=["Health"])
async def readiness_check() -> ReadyResponse:
    """
    Readiness probe endpoint.

    Returns if the service is ready to accept requests.
    Used by Kubernetes readiness probes.
    """
    return ReadyResponse()


@app.post(
    "/extract",
    response_model=ExtractResponse,
    tags=["Extraction"],
    status_code=status.HTTP_200_OK,
    responses={
        400: {"model": ErrorDetail, "description": "Invalid file type or corrupted file"},
        500: {"model": ErrorDetail, "description": "Internal extraction error"},
        502: {"model": ErrorDetail, "description": "Failed to download file"},
    },
)
async def extract_text(request: ExtractRequest) -> ExtractResponse:
    """
    Extract text from a file.

    Downloads the file from the provided URL and extracts text content
    based on the file type (PDF, DOCX, or PPTX).
    
    Returns extracted text along with metadata including page/slide count
    and word count.
    """
    file_path: Path | None = None
    
    try:
        # Map file types to extensions
        extension_map = {
            "pdf": ".pdf",
            "docx": ".docx",
            "pptx": ".pptx",
        }
        
        # Download the file
        try:
            file_path = await download_file(
                str(request.file_url), 
                extension_map[request.file_type]
            )
        except httpx.HTTPStatusError as e:
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail=f"Failed to download file: HTTP {e.response.status_code}",
            )
        except httpx.RequestError as e:
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail=f"Failed to download file: {str(e)}",
            )
        
        # Extract text based on file type
        try:
            if request.file_type == "pdf":
                text, pages = extract_pdf(file_path)
            elif request.file_type == "docx":
                text, pages = extract_docx(file_path)
            elif request.file_type == "pptx":
                text, pages = extract_pptx(file_path)
            else:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Unsupported file type: {request.file_type}",
                )
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Failed to extract text from {request.file_type}: {str(e)}",
            )
        
        # Calculate metadata
        word_count = count_words(text)
        char_count = len(text)
        
        return ExtractResponse(
            text=text,
            metadata=ExtractMetadata(
                pages=pages,
                word_count=word_count,
                char_count=char_count,
            ),
        )
    
    finally:
        # Clean up temporary file
        if file_path and file_path.exists():
            file_path.unlink()


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
